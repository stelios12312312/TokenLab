"""
Z1 M2 Market Dynamics — PPTX Presentation Generator

Generates a polished slide deck from the M2 simulation results.
Usage:
    PYTHONPATH=src:. python -m projects.z1.m2_market_dynamics.generate_pptx [--run-dir outputs/z1_core_solvency/<run_id>]
"""

import argparse
import json
import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colours ────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)
SURFACE   = RGBColor(0x1E, 0x29, 0x3B)
ACCENT    = RGBColor(0x3B, 0x82, 0xF6)
TEXT_MAIN = RGBColor(0xE2, 0xE8, 0xF0)
MUTED     = RGBColor(0x94, 0xA3, 0xB8)
SAFE      = RGBColor(0x16, 0xA3, 0x4A)
WARN      = RGBColor(0xF5, 0x9E, 0x0B)
DANGER    = RGBColor(0xDC, 0x26, 0x26)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_slide_bg(slide, color=BG_DARK):
    """Set solid background colour on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_notes(slide, text):
    """Set speaker notes on a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=TEXT_MAIN, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Helper: add a styled textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_bullet_list(slide, left, top, width, height, items,
                     font_size=16, color=TEXT_MAIN):
    """Helper: add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return tf


def _add_table(slide, left, top, width, height, headers, rows,
               header_color=ACCENT, row_color=TEXT_MAIN):
    """Helper: add a styled table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Style header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x25, 0x33, 0x48)

    # Style rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = row_color
                paragraph.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE if i % 2 == 0 else BG_DARK

    return table


def _fmt(val, fmt_str='.2f'):
    if val is None:
        return 'N/A'
    try:
        return f'{val:{fmt_str}}'
    except (ValueError, TypeError):
        return str(val)


def generate_pptx(run_dir: str, output_path: str = None):
    """Generate the M2 PPTX from a run directory."""

    # Load summaries
    summaries = {}
    for jf in glob.glob(os.path.join(run_dir, "*_summary.json")):
        name = os.path.basename(jf).replace("_summary.json", "")
        with open(jf) as f:
            summaries[name] = json.load(f)

    if not summaries:
        raise FileNotFoundError(f"No summary JSONs found in {run_dir}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                 "Z1 M2 Market Dynamics", font_size=44, bold=True, color=WHITE)
    _add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                 "Endogenous Pricing & Adversarial Solvency Model", font_size=24, color=ACCENT)
    _add_textbox(slide, Inches(1), Inches(5), Inches(11), Inches(0.5),
                 "TokenLab Simulation Report", font_size=16, color=MUTED)
    _set_notes(slide,
        "M2 is the second generation of our solvency model. "
        "M1 tested pure structural integrity with static assumptions. "
        "M2 introduces endogenous market pricing via an AMM, adversarial agent behavior (panic selling), "
        "and a campaign escrow engine. The goal is to answer: can the system survive when price is determined "
        "by the market and agents act against the protocol's interest?")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 2: What M2 Models
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "What M2 Models", font_size=32, bold=True, color=WHITE)
    _add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
                 "Core loop: Issuance → Vesting → Settlement → AMM (Price) → Panic Feedback → Utility Recap",
                 font_size=14, color=MUTED)

    _add_table(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4),
               ["Component", "Implementation"],
               [
                   ["AMM Pricing", "Endogenous price discovery via constant-product Z1U/USD pool"],
                   ["Escrow Engine", "Brand deposits go to escrow; 25% fee to Treasury"],
                   ["Panic Triggers", "10% price drop triggers 'Bank Run' (10x settlement surge)"],
                   ["L6 Floor Guard", "Constitutional 25% AR floor enforced by settlement capping"],
                   ["Utility Release", "Escrowed Z1U released to Treasury upon user utility spend"],
                   ["Health Throttle", "Reduces ACR issuance when AR ratio < 0.3"],
                   ["Invariants", "Full conservation including AMM reserves and Escrow balances"],
               ])
    _set_notes(slide,
        "Key upgrade from M1: price is no longer an input — it emerges from a constant-product AMM pool. "
        "When users settle their ACR rewards for Z1U and sell on the AMM, the price drops. "
        "If the price drops more than 10%, a panic trigger fires and settlement demand surges 10x. "
        "The L6 floor guard is our constitutional defense: settlement is dynamically capped so the "
        "Audience Reserve never falls below 25% of circulating supply. "
        "The escrow engine means brand deposits don't go directly to users — 25% is captured as a Treasury fee upfront.")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 3: Scenario Results (one per scenario)
    # ═══════════════════════════════════════════════════════════════════
    for name, summary in summaries.items():
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide)

        classification = summary.get('classification', 'N/A')
        cls_color = {'collapse': DANGER, 'stressed': WARN, 'stable': SAFE}.get(classification, MUTED)

        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
                     f"Scenario: {name.replace('_', ' ').title()}", font_size=32, bold=True, color=WHITE)
        _add_textbox(slide, Inches(9), Inches(0.4), Inches(4), Inches(0.8),
                     classification.upper(), font_size=28, bold=True, color=cls_color, alignment=PP_ALIGN.RIGHT)

        _add_table(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5),
                   ["Metric", "Value"],
                   [
                       ["Final AR Ratio", _fmt(summary.get('final_ar_ratio'))],
                       ["Min AR Ratio", _fmt(summary.get('min_ar_ratio'))],
                       ["Final Z1U Price", "$" + _fmt(summary.get('final_price'), '.4f')],
                       ["Min Z1U Price", "$" + _fmt(summary.get('min_price'), '.4f')],
                       ["Final Treasury", _fmt(summary.get('final_treasury'), ',.0f') + " Z1U"],
                       ["Final Escrow", _fmt(summary.get('final_escrow'), ',.0f') + " Z1U"],
                       ["Max Settlement Queue", _fmt(summary.get('max_settlement_queue_z1u'), ',.0f') + " Z1U"],
                       ["Total Utility Spend", _fmt(summary.get('total_utility_spend'), ',.0f') + " Z1U"],
                       ["Total Burn", _fmt(summary.get('total_burn'), ',.0f') + " Z1U"],
                       ["Total Brand Deposits", _fmt(summary.get('total_brand_inflow'), ',.0f') + " Z1U"],
                       ["Throttle Epochs", f"{summary.get('throttle_epochs', 0)} / 104"],
                       ["Panic Epochs", f"{summary.get('panic_epochs', 0)} / 104"],
                   ])

        # Add plot images if available
        plot_dirs = [
            os.path.join(run_dir, "plots", name),
            os.path.join(run_dir, "plots"),
        ]
        for pd_path in plot_dirs:
            if os.path.isdir(pd_path):
                pngs = sorted(glob.glob(os.path.join(pd_path, "*.png")))[:2]
                for i, png in enumerate(pngs):
                    try:
                        slide.shapes.add_picture(
                            png,
                            Inches(6.8) + Inches(i * 0.1),
                            Inches(1.5),
                            width=Inches(5.5),
                        )
                    except Exception:
                        pass
                break

        # Speaker notes for scenario slides
        if classification == 'collapse':
            notes = (
                f"The {name} scenario demonstrates a structural failure. "
                f"The AR ratio fell to {_fmt(summary.get('min_ar_ratio'))}, below the collapse threshold. "
                f"The throttle engaged for {summary.get('throttle_epochs', 0)} of 104 epochs but couldn't save it. "
                f"Key takeaway: the inflow was insufficient to sustain the outflow pressure."
            )
        elif classification == 'stressed':
            notes = (
                f"The {name} scenario shows the system under heavy pressure but surviving. "
                f"AR ratio bottomed at {_fmt(summary.get('min_ar_ratio'))} and recovered to {_fmt(summary.get('final_ar_ratio'))}. "
                f"The throttle was active for {summary.get('throttle_epochs', 0)} of 104 epochs — this is the system's immune response. "
                f"Settlement queue peaked at {_fmt(summary.get('max_settlement_queue_z1u'), ',.0f')} Z1U. "
                f"The Treasury went to zero but the AR floor held, proving the constitutional guard works."
            )
        else:
            notes = (
                f"The {name} scenario shows healthy loop dynamics. "
                f"Utility spend, brand inflows, and Treasury fees sustain the AR at {_fmt(summary.get('final_ar_ratio'))}. "
                f"No throttle was needed. This is the target operating state."
            )
        _set_notes(slide, notes)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE: Comparison (if multiple scenarios)
    # ═══════════════════════════════════════════════════════════════════
    if len(summaries) > 1:
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide)
        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                     "Scenario Comparison", font_size=32, bold=True, color=WHITE)

        headers = ["Metric"] + [n.replace('_', ' ').title() for n in summaries.keys()]
        metrics_to_compare = [
            ("Final AR Ratio", 'final_ar_ratio', '.2f'),
            ("Min AR Ratio", 'min_ar_ratio', '.2f'),
            ("Final Treasury", 'final_treasury', ',.0f'),
            ("Max Queue (Z1U)", 'max_settlement_queue_z1u', ',.0f'),
            ("Total Brand Deposits", 'total_brand_inflow', ',.0f'),
            ("Throttle Epochs", 'throttle_epochs', '.0f'),
            ("Panic Epochs", 'panic_epochs', '.0f'),
            ("Classification", 'classification', None),
        ]
        rows = []
        for label, key, fmt_str in metrics_to_compare:
            row = [label]
            for s in summaries.values():
                val = s.get(key, 'N/A')
                if fmt_str and val != 'N/A':
                    row.append(_fmt(val, fmt_str))
                else:
                    row.append(str(val).upper() if key == 'classification' else str(val))
            rows.append(row)

        _add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5),
                   headers, rows)
        _set_notes(slide,
            "This side-by-side comparison reveals the system's resilience envelope. "
            "Notice the bank run scenario has 50x less brand inflow but the AR only drops from 1.17 to 0.88. "
            "The settlement queue is 3x larger under bank run conditions, yet the system holds. "
            "This proves the throttle-and-cap mechanism works. "
            "The key insight: it's not the size of the attack that matters, "
            "it's whether the constitutional floor is mechanically enforced.")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE: Top Solvency Drivers
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Top Solvency Drivers (OAT Elasticity)", font_size=32, bold=True, color=WHITE)
    _add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
                 "Only 3 of 12 parameters have meaningful elasticity. Everything else is noise.",
                 font_size=14, color=MUTED)

    _add_table(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(3),
               ["Rank", "Parameter", "AR Elasticity", "Meaning"],
               [
                   ["1", "treasury_topup_threshold_ratio", "+2.98", "Most powerful lever — controls when Treasury recapitalises AR"],
                   ["2", "audience_reserve_initial", "−0.77", "Bigger starting AR makes topups less effective (denominator effect)"],
                   ["3", "brand_inflow_per_epoch", "+0.42", "External revenue — the oxygen of the system"],
                   ["4–12", "All others", "< 0.05", "Settlement cap, vesting, fees, burn, throttle — near-zero impact"],
               ])
    _set_notes(slide,
        "We ran One-At-a-Time sensitivity screening across 12 parameters, perturbing each by ±20%. "
        "The result is striking: only 3 parameters materially affect the final AR ratio. "
        "The topup threshold has 3x the elasticity of brand inflow — meaning WHEN you start recapitalising "
        "the AR matters far more than how much money flows in. "
        "The counterintuitive finding: a bigger starting AR actually hurts the ratio because it inflates "
        "the denominator, making each topup proportionally smaller. "
        "Everything else — settlement caps, vesting periods, fee percentages, burn rates — is noise.")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE: Hard Constraints
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Hard Constraints — Violate = Collapse", font_size=32, bold=True, color=DANGER)

    _add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.5),
               ["Lock", "Rule", "Minimum Threshold"],
               [
                   ["L1", "Solvency Ratio (outflow / inflow)", "< 0.8"],
                   ["L3", "Brand Inflow / AR per epoch", "≥ 1%"],
                   ["L6", "AR / Circulating Supply", "≥ 25%"],
                   ["L9", "Max single-epoch AR drain", "≤ 10% of initial AR"],
               ])

    _add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
                 "Optimal Parameters (Monte Carlo Calibration)", font_size=22, bold=True, color=WHITE)

    _add_table(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(2),
               ["Parameter", "Optimal Value", "Safe Range"],
               [
                   ["Solvency Ratio", "0.006", "< 0.8"],
                   ["Settlement Ratio", "0.10", "0.05 – 0.75"],
                   ["Utility Fee Share", "34%", "6% – 40%"],
                   ["Brand Inflow", "2.24% of AR", "≥ 1% (absolute minimum)"],
                   ["Campaign Fee", "25%", "≥ 15% for AR floor defense"],
               ])
    _set_notes(slide,
        "These are the hard lines from the simulation. L1 is the master invariant: if outflow/inflow exceeds 0.8, "
        "you're in fragile territory. Above 1.0, collapse is guaranteed in every observed case. "
        "L3 is the survival floor: below 1% brand inflow per epoch relative to AR, not a single scenario survived. "
        "L6 is constitutional: the 25% AR floor must be enforced by code, not by governance vote. "
        "The optimal calibration from our Monte Carlo search shows the safe operating ranges. "
        "Note the settlement ratio optimum is 0.10 — much lower than the 0.5 default many assume.")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE: Design Rules
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Design Rules for Z1 Tokenomics", font_size=32, bold=True, color=WHITE)

    _add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5), [
        "1. Set the topup trigger aggressively — sensitivity to WHEN you recapitalise is 3x higher than to inflow amount",
        "2. Don't over-capitalise the AR at launch — a giant starting reserve creates a false sense of security",
        "3. Brand inflow is the oxygen — below 1% of AR per epoch, the system collapses in every observed case",
        "4. Constitutional 25% AR floor must be enforced mechanically — don't trust governance to maintain it",
        "5. Monitor passive viewer cohort share — if extractors grow, raise their settlement friction",
    ], font_size=18)
    _set_notes(slide,
        "These five rules are the practical output of the M2 simulation. "
        "Rule 1 is the most important and most counterintuitive: it's not about how much money you have, "
        "it's about how early you trigger the recapitalisation mechanism. A 1% change in the trigger threshold "
        "produces a 3% change in the final AR ratio. "
        "Rule 2 challenges the instinct to launch with a massive reserve — the math works against you. "
        "Rule 3 is non-negotiable: without external revenue, the system is a Ponzi by construction. "
        "Rule 4 is about trust minimisation: if the floor can be changed by governance, it will be under attack. "
        "Rule 5 flags the cohort composition risk — growth strategies that attract extractors are system-hostile.")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE: Passive Viewer Problem
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "The Passive Viewer Problem", font_size=32, bold=True, color=WARN)

    _add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2),
               ["Cohort", "Settle Propensity", "Spend Rate", "Settle/Spend Ratio", "Status"],
               [
                   ["Passive Viewers", "0.25", "0.10", "2.50x", "⚠ Net Extractor"],
                   ["Active Viewers", "0.30", "0.40", "0.75x", "⚠ Net Extractor"],
                   ["Power Users", "0.15", "0.80", "0.19x", "✓ Net Contributor"],
               ])

    _add_textbox(slide, Inches(0.8), Inches(4), Inches(11), Inches(0.5),
                 "Target: settle/spend ≤ 0.5x per cohort", font_size=16, color=MUTED)

    _add_bullet_list(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(2), [
        "System survives because power users subsidise passive viewers",
        "Brand inflow covers the remaining gap",
        "If passive viewer share grows (e.g. referral campaigns), the system breaks",
    ], font_size=16, color=TEXT_MAIN)
    _set_notes(slide,
        "This is the uncomfortable finding. Passive viewers extract 2.5x more value than they contribute. "
        "The L4 target says settle/spend should be below 0.5x — passive viewers are at 2.5x, five times over. "
        "The system only works because: (a) power users are massive net contributors at 0.19x, "
        "(b) brand inflow subsidises the gap, and (c) the 60/30/10 population split means passive viewers "
        "are individually small. But any growth strategy that shifts this mix — referral bonuses, airdrops, "
        "gamification that attracts speculators — is structurally dangerous. "
        "The recommendation: tie settlement friction to cohort behaviour, not just global caps.")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE: M3 Roadmap
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Recommended M3 Extensions", font_size=32, bold=True, color=WHITE)

    _add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(3.5),
               ["Priority", "Extension", "Rationale"],
               [
                   ["P0", "Multi-token Staking (Z1P)", "Sink for Z1U to reduce immediate sell pressure"],
                   ["P0", "Dynamic LP Dynamics", "Allow AMM liquidity to respond to yield and volatility"],
                   ["P1", "Full PCS Decomposition", "Replace reduced-form verification with multi-factor scoring"],
                   ["P1", "Creator & Validator Cohorts", "Expand economy to include supply-side agents"],
                   ["P2", "Governance & Delegation", "Model attack vectors on Treasury and Protocol parameters"],
                   ["P3", "Prediction Market Sink", "Integrate secondary utility sinks for Z1U"],
               ])
    _set_notes(slide,
        "M3 priorities are driven by the gaps M2 exposed. "
        "P0 items address the biggest structural weakness: there's no sink for Z1U sell pressure other than utility spend. "
        "Z1P staking would give users a reason to hold rather than sell. Dynamic LP would make the AMM more realistic. "
        "P1 items expand the agent taxonomy — creators and validators are currently missing from the economy. "
        "P2 governance modeling is critical because the constitutional floor (L6) is only as strong as the "
        "governance mechanism that protects it. If an attacker can vote to lower the floor, the system is vulnerable. "
        "P3 prediction markets are a secondary utility sink — nice to have, not essential.")

    # ═══════════════════════════════════════════════════════════════════
    # Save
    # ═══════════════════════════════════════════════════════════════════
    if output_path is None:
        output_path = os.path.join(run_dir, "Z1_M2_Market_Dynamics.pptx")

    prs.save(output_path)
    print(f"PPTX saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate Z1 M2 PPTX presentation")
    parser.add_argument("--run-dir", type=str, default=None,
                        help="Path to simulation run directory (auto-detects latest if omitted)")
    parser.add_argument("--output", type=str, default=None,
                        help="Output PPTX path (defaults to <run-dir>/Z1_M2_Market_Dynamics.pptx)")
    args = parser.parse_args()

    if args.run_dir is None:
        # Auto-detect latest run
        base = os.path.join("outputs", "z1_core_solvency")
        runs = sorted(glob.glob(os.path.join(base, "*")))
        if not runs:
            raise FileNotFoundError(f"No run directories found in {base}")
        args.run_dir = runs[-1]
        print(f"Auto-detected latest run: {args.run_dir}")

    generate_pptx(args.run_dir, args.output)
