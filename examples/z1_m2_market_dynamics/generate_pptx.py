"""
Z1 M2 Market Dynamics — PPTX Presentation Generator

Generates a polished slide deck from the M2 simulation results.
Usage:
    PYTHONPATH=src:. python -m examples.z1_m2_market_dynamics.generate_pptx [--run-dir outputs/z1_core_solvency/<run_id>]
"""

import argparse
import json
import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colours ────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)
SURFACE   = RGBColor(0x1E, 0x29, 0x3B)
ACCENT    = RGBColor(0x3B, 0x82, 0xF6)
TEXT_MAIN = RGBColor(0xE2, 0xE8, 0xF0)
MUTED     = RGBColor(0x94, 0xA3, 0xB8)
SAFE      = RGBColor(0x16, 0xA3, 0x4A)
WARN      = RGBColor(0xF5, 0x9E, 0x0B)
DANGER    = RGBColor(0xDC, 0x26, 0x26)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_slide_bg(slide, color=BG_DARK):
    """Set solid background colour on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=TEXT_MAIN, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Helper: add a styled textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_bullet_list(slide, left, top, width, height, items,
                     font_size=16, color=TEXT_MAIN):
    """Helper: add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return tf


def _add_table(slide, left, top, width, height, headers, rows,
               header_color=ACCENT, row_color=TEXT_MAIN):
    """Helper: add a styled table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Style header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x25, 0x33, 0x48)

    # Style rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = row_color
                paragraph.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE if i % 2 == 0 else BG_DARK

    return table


def _fmt(val, fmt_str='.2f'):
    if val is None:
        return 'N/A'
    try:
        return f'{val:{fmt_str}}'
    except (ValueError, TypeError):
        return str(val)


def generate_pptx(run_dir: str, output_path: str = None):
    """Generate the M2 PPTX from a run directory."""

    # Load summaries
    summaries = {}
    for jf in glob.glob(os.path.join(run_dir, "*_summary.json")):
        name = os.path.basename(jf).replace("_summary.json", "")
        with open(jf) as f:
            summaries[name] = json.load(f)

    if not summaries:
        raise FileNotFoundError(f"No summary JSONs found in {run_dir}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                 "Z1 M2 Market Dynamics", font_size=44, bold=True, color=WHITE)
    _add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                 "Endogenous Pricing & Adversarial Solvency Model", font_size=24, color=ACCENT)
    _add_textbox(slide, Inches(1), Inches(5), Inches(11), Inches(0.5),
                 "TokenLab Simulation Report", font_size=16, color=MUTED)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 2: What M2 Models
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "What M2 Models", font_size=32, bold=True, color=WHITE)
    _add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
                 "Core loop: Issuance → Vesting → Settlement → AMM (Price) → Panic Feedback → Utility Recap",
                 font_size=14, color=MUTED)

    _add_table(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4),
               ["Component", "Implementation"],
               [
                   ["AMM Pricing", "Endogenous price discovery via constant-product Z1U/USD pool"],
                   ["Escrow Engine", "Brand deposits go to escrow; 25% fee to Treasury"],
                   ["Panic Triggers", "10% price drop triggers 'Bank Run' (10x settlement surge)"],
                   ["L6 Floor Guard", "Constitutional 25% AR floor enforced by settlement capping"],
                   ["Utility Release", "Escrowed Z1U released to Treasury upon user utility spend"],
                   ["Health Throttle", "Reduces ACR issuance when AR ratio < 0.3"],
                   ["Invariants", "Full conservation including AMM reserves and Escrow balances"],
               ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 3: Scenario Results (one per scenario)
    # ═══════════════════════════════════════════════════════════════════
    for name, summary in summaries.items():
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide)

        classification = summary.get('classification', 'N/A')
        cls_color = {'collapse': DANGER, 'stressed': WARN, 'stable': SAFE}.get(classification, MUTED)

        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
                     f"Scenario: {name.replace('_', ' ').title()}", font_size=32, bold=True, color=WHITE)
        _add_textbox(slide, Inches(9), Inches(0.4), Inches(4), Inches(0.8),
                     classification.upper(), font_size=28, bold=True, color=cls_color, alignment=PP_ALIGN.RIGHT)

        _add_table(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5),
                   ["Metric", "Value"],
                   [
                       ["Final AR Ratio", _fmt(summary.get('final_ar_ratio'))],
                       ["Min AR Ratio", _fmt(summary.get('min_ar_ratio'))],
                       ["Final Z1U Price", "$" + _fmt(summary.get('final_price'), '.4f')],
                       ["Min Z1U Price", "$" + _fmt(summary.get('min_price'), '.4f')],
                       ["Final Treasury", _fmt(summary.get('final_treasury'), ',.0f') + " Z1U"],
                       ["Final Escrow", _fmt(summary.get('final_escrow'), ',.0f') + " Z1U"],
                       ["Max Settlement Queue", _fmt(summary.get('max_settlement_queue_z1u'), ',.0f') + " Z1U"],
                       ["Total Utility Spend", _fmt(summary.get('total_utility_spend'), ',.0f') + " Z1U"],
                       ["Total Burn", _fmt(summary.get('total_burn'), ',.0f') + " Z1U"],
                       ["Total Brand Deposits", _fmt(summary.get('total_brand_inflow'), ',.0f') + " Z1U"],
                       ["Throttle Epochs", f"{summary.get('throttle_epochs', 0)} / 104"],
                       ["Panic Epochs", f"{summary.get('panic_epochs', 0)} / 104"],
                   ])

        # Add plot images if available
        plot_dirs = [
            os.path.join(run_dir, "plots", name),
            os.path.join(run_dir, "plots"),
        ]
        for pd_path in plot_dirs:
            if os.path.isdir(pd_path):
                pngs = sorted(glob.glob(os.path.join(pd_path, "*.png")))[:2]
                for i, png in enumerate(pngs):
                    try:
                        slide.shapes.add_picture(
                            png,
                            Inches(6.8) + Inches(i * 0.1),
                            Inches(1.5),
                            width=Inches(5.5),
                        )
                    except Exception:
                        pass
                break

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE: Comparison (if multiple scenarios)
    # ═══════════════════════════════════════════════════════════════════
    if len(summaries) > 1:
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide)
        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                     "Scenario Comparison", font_size=32, bold=True, color=WHITE)

        headers = ["Metric"] + [n.replace('_', ' ').title() for n in summaries.keys()]
        metrics_to_compare = [
            ("Final AR Ratio", 'final_ar_ratio', '.2f'),
            ("Min AR Ratio", 'min_ar_ratio', '.2f'),
            ("Final Treasury", 'final_treasury', ',.0f'),
            ("Max Queue (Z1U)", 'max_settlement_queue_z1u', ',.0f'),
            ("Total Brand Deposits", 'total_brand_inflow', ',.0f'),
            ("Throttle Epochs", 'throttle_epochs', '.0f'),
            ("Panic Epochs", 'panic_epochs', '.0f'),
            ("Classification", 'classification', None),
        ]
        rows = []
        for label, key, fmt_str in metrics_to_compare:
            row = [label]
            for s in summaries.values():
                val = s.get(key, 'N/A')
                if fmt_str and val != 'N/A':
                    row.append(_fmt(val, fmt_str))
                else:
                    row.append(str(val).upper() if key == 'classification' else str(val))
            rows.append(row)

        _add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5),
                   headers, rows)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE: Top Solvency Drivers
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Top Solvency Drivers (OAT Elasticity)", font_size=32, bold=True, color=WHITE)
    _add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
                 "Only 3 of 12 parameters have meaningful elasticity. Everything else is noise.",
                 font_size=14, color=MUTED)

    _add_table(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(3),
               ["Rank", "Parameter", "AR Elasticity", "Meaning"],
               [
                   ["1", "treasury_topup_threshold_ratio", "+2.98", "Most powerful lever — controls when Treasury recapitalises AR"],
                   ["2", "audience_reserve_initial", "−0.77", "Bigger starting AR makes topups less effective (denominator effect)"],
                   ["3", "brand_inflow_per_epoch", "+0.42", "External revenue — the oxygen of the system"],
                   ["4–12", "All others", "< 0.05", "Settlement cap, vesting, fees, burn, throttle — near-zero impact"],
               ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE: Hard Constraints
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Hard Constraints — Violate = Collapse", font_size=32, bold=True, color=DANGER)

    _add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.5),
               ["Lock", "Rule", "Minimum Threshold"],
               [
                   ["L1", "Solvency Ratio (outflow / inflow)", "< 0.8"],
                   ["L3", "Brand Inflow / AR per epoch", "≥ 1%"],
                   ["L6", "AR / Circulating Supply", "≥ 25%"],
                   ["L9", "Max single-epoch AR drain", "≤ 10% of initial AR"],
               ])

    _add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
                 "Optimal Parameters (Monte Carlo Calibration)", font_size=22, bold=True, color=WHITE)

    _add_table(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(2),
               ["Parameter", "Optimal Value", "Safe Range"],
               [
                   ["Solvency Ratio", "0.006", "< 0.8"],
                   ["Settlement Ratio", "0.10", "0.05 – 0.75"],
                   ["Utility Fee Share", "34%", "6% – 40%"],
                   ["Brand Inflow", "2.24% of AR", "≥ 1% (absolute minimum)"],
                   ["Campaign Fee", "25%", "≥ 15% for AR floor defense"],
               ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE: Design Rules
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Design Rules for Z1 Tokenomics", font_size=32, bold=True, color=WHITE)

    _add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5), [
        "1. Set the topup trigger aggressively — sensitivity to WHEN you recapitalise is 3x higher than to inflow amount",
        "2. Don't over-capitalise the AR at launch — a giant starting reserve creates a false sense of security",
        "3. Brand inflow is the oxygen — below 1% of AR per epoch, the system collapses in every observed case",
        "4. Constitutional 25% AR floor must be enforced mechanically — don't trust governance to maintain it",
        "5. Monitor passive viewer cohort share — if extractors grow, raise their settlement friction",
    ], font_size=18)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE: Passive Viewer Problem
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "The Passive Viewer Problem", font_size=32, bold=True, color=WARN)

    _add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2),
               ["Cohort", "Settle Propensity", "Spend Rate", "Settle/Spend Ratio", "Status"],
               [
                   ["Passive Viewers", "0.25", "0.10", "2.50x", "⚠ Net Extractor"],
                   ["Active Viewers", "0.30", "0.40", "0.75x", "⚠ Net Extractor"],
                   ["Power Users", "0.15", "0.80", "0.19x", "✓ Net Contributor"],
               ])

    _add_textbox(slide, Inches(0.8), Inches(4), Inches(11), Inches(0.5),
                 "Target: settle/spend ≤ 0.5x per cohort", font_size=16, color=MUTED)

    _add_bullet_list(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(2), [
        "System survives because power users subsidise passive viewers",
        "Brand inflow covers the remaining gap",
        "If passive viewer share grows (e.g. referral campaigns), the system breaks",
    ], font_size=16, color=TEXT_MAIN)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE: M3 Roadmap
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Recommended M3 Extensions", font_size=32, bold=True, color=WHITE)

    _add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(3.5),
               ["Priority", "Extension", "Rationale"],
               [
                   ["P0", "Multi-token Staking (Z1P)", "Sink for Z1U to reduce immediate sell pressure"],
                   ["P0", "Dynamic LP Dynamics", "Allow AMM liquidity to respond to yield and volatility"],
                   ["P1", "Full PCS Decomposition", "Replace reduced-form verification with multi-factor scoring"],
                   ["P1", "Creator & Validator Cohorts", "Expand economy to include supply-side agents"],
                   ["P2", "Governance & Delegation", "Model attack vectors on Treasury and Protocol parameters"],
                   ["P3", "Prediction Market Sink", "Integrate secondary utility sinks for Z1U"],
               ])

    # ═══════════════════════════════════════════════════════════════════
    # Save
    # ═══════════════════════════════════════════════════════════════════
    if output_path is None:
        output_path = os.path.join(run_dir, "Z1_M2_Market_Dynamics.pptx")

    prs.save(output_path)
    print(f"PPTX saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate Z1 M2 PPTX presentation")
    parser.add_argument("--run-dir", type=str, default=None,
                        help="Path to simulation run directory (auto-detects latest if omitted)")
    parser.add_argument("--output", type=str, default=None,
                        help="Output PPTX path (defaults to <run-dir>/Z1_M2_Market_Dynamics.pptx)")
    args = parser.parse_args()

    if args.run_dir is None:
        # Auto-detect latest run
        base = os.path.join("outputs", "z1_core_solvency")
        runs = sorted(glob.glob(os.path.join(base, "*")))
        if not runs:
            raise FileNotFoundError(f"No run directories found in {base}")
        args.run_dir = runs[-1]
        print(f"Auto-detected latest run: {args.run_dir}")

    generate_pptx(args.run_dir, args.output)
